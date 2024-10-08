import streamlit as st
import hashlib
import json
import os
import re
import openai
import pandas as pd
import pptx
import PyPDF2
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
import time
from openai import OpenAI

# Load environment variables from .env file
load_dotenv()

# Read API keys from environment variables
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Initialize OpenAI
openai.api_key = OPENAI_API_KEY
client = OpenAI()

# Initialize Pinecone instance
pc = Pinecone(api_key=PINECONE_API_KEY)

# Function to hash passwords
def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

# Load persistent storage
def load_storage():
    if os.path.exists("storage.json"):
        with open("storage.json", "r") as f:
            return json.load(f)
    return {'users': {}, 'collections': {}, 'uploaded_files': {}}

# Save persistent storage
def save_storage(storage):
    with open("storage.json", "w") as f:
        json.dump(storage, f, indent=4)

# Initialize persistent storage
storage = load_storage()

# Initialize session state
if 'users' not in st.session_state:
    st.session_state.users = storage['users']
if 'current_user' not in st.session_state:
    st.session_state.current_user = None  # To track the logged-in user
if 'collections' not in st.session_state:
    st.session_state.collections = storage['collections']
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = storage['uploaded_files']
if 'selected_collection' not in st.session_state:
    st.session_state.selected_collection = None  # Track selected collection
if 'page' not in st.session_state:
    st.session_state.page = 'login'  # Default to the login page
if 'index_name' not in st.session_state:
    st.session_state.index_name = None  # Store the Pinecone index name input by the user

# User registration
def register(username, password):
    if username in st.session_state.users:
        return False  # Username already exists
    st.session_state.users[username] = hash_password(password)
    st.session_state.collections[username] = []
    st.session_state.uploaded_files[username] = {}
    save_storage({'users': st.session_state.users, 'collections': st.session_state.collections, 'uploaded_files': st.session_state.uploaded_files})
    return True

# User login
def login(username, password):
    if username in st.session_state.users and st.session_state.users[username] == hash_password(password):
        st.session_state.current_user = username
        return True
    return False

# Validate index name
def validate_index_name(index_name):
    # Check if the index name conforms to Pinecone's naming rules
    pattern = r'^[a-z0-9-]+$'
    return bool(re.match(pattern, index_name))

# Add collection
def add_collection(username, collection_name, index_name, uploaded_files):
    collections = st.session_state.collections.get(username, [])

    # Ensure the collection name and index name pair is unique
    for collection in collections:
        if collection['name'] == collection_name:
            return "Collection name already exists."
        if collection['index_name'] == index_name:
            return "Index name already exists."

    # Attempt to create the Pinecone index and upsert files
    result = process_and_upsert_files(username, collection_name, index_name, uploaded_files)
    if result != "Success":
        return result  # Return the error message if index creation/upsert failed

    # Only add the collection if Pinecone index creation and upsert are successful
    collections.append({'name': collection_name, 'index_name': index_name})
    st.session_state.collections[username] = collections
    st.session_state.uploaded_files[username][collection_name] = [file.name for file in uploaded_files]
    
    # Debugging print statements to verify storage content before saving
    print("Debug: Before saving storage.json")
    print("Users:", st.session_state.users)
    print("Collections:", st.session_state.collections)
    print("Uploaded Files:", st.session_state.uploaded_files)
    
    save_storage({
        'users': st.session_state.users,
        'collections': st.session_state.collections,
        'uploaded_files': st.session_state.uploaded_files
    })

    # Debugging print statements to verify that data was saved
    print("Debug: After saving storage.json")
    return "Success"

# Function to process and upsert files into Pinecone
def process_and_upsert_files(username, collection_name, index_name, uploaded_files):
    try:
        if index_name not in pc.list_indexes().names():
            pc.create_index(
                name=index_name,
                dimension=1536,  # Replace with appropriate dimension
                metric='cosine',  # or another metric such as 'cosine'
                spec=ServerlessSpec(cloud='aws', region='us-east-1')
            )
        while not pc.describe_index(index_name).status['ready']:
            time.sleep(1)
        st.session_state.index = pc.Index(index_name)

        for uploaded_file in uploaded_files:
            if uploaded_file.type == "text/plain":
                process_txt_file(uploaded_file, collection_name)
            elif uploaded_file.type == "application/pdf":
                process_pdf_file(uploaded_file, collection_name)
            elif uploaded_file.type == "text/csv":
                process_csv_file(uploaded_file, collection_name)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                process_ppt_file(uploaded_file, collection_name)
            else:
                st.warning(f"Unsupported file type: {uploaded_file.type}")

        return "Success"
    except Exception as e:
        st.error(f"An error occurred: {e}")
        return f"Failed to create index or upsert files: {e}"

# Function to process and upsert a text file
def process_txt_file(file, collection_name):
    text = file.read().decode("utf-8")
    upsert_to_pinecone(text, collection_name)

# Function to process and upsert a PDF file
def process_pdf_file(file, collection_name):
    reader = PyPDF2.PdfFileReader(file)
    text = ""
    for page_num in range(reader.getNumPages()):
        page = reader.getPage(page_num)
        
        # Attempt to use the appropriate method based on the version of PyPDF2
        try:
            # For newer versions
            text += page.extract_text()
        except AttributeError:
            # For older versions
            text += page.extractText()
    upsert_to_pinecone(text, collection_name)

# Function to process and upsert a CSV file
def process_csv_file(file, collection_name):
    df = pd.read_csv(file)
    for _, row in df.iterrows():
        text = " ".join(row.astype(str).tolist())
        upsert_to_pinecone(text, collection_name)

# Function to process and upsert a PPT file
def process_ppt_file(file, collection_name):
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    upsert_to_pinecone(text, collection_name)

# Function to upsert text into Pinecone
def upsert_to_pinecone(text, collection_name):
    chunks = [text[i:i + 2000] for i in range(0, len(text), 2000)]
    for chunk in chunks:
        # Create an embedding for the chunk using the client and response handling
        chunk = chunk.replace("\n", " ")  # Ensure no newlines in the text
        response = client.embeddings.create(input=[chunk], model="text-embedding-ada-002")
       
        # Extract the embedding from the response
        embedding = response.data[0].embedding
        
        # Create a unique ID for the chunk using SHA-256
        id = hashlib.sha256(chunk.encode()).hexdigest()
        
        # Upsert the embedding into the Pinecone index
        st.session_state.index.upsert(vectors=[{"id": id,"values": embedding,"metadata": {"text": chunk}}])

# Pages
def show_login_page():
    st.title("Login")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")

    if st.button("Login"):
        if login(username, password):
            st.success("Logged in successfully")
            st.session_state.page = 'collection_management'
        else:
            st.error("Invalid username or password")

    if st.button("Go to Register"):
        st.session_state.page = 'register'

def show_register_page():
    st.title("Register")
    new_username = st.text_input("New Username")
    new_password = st.text_input("New Password", type="password")

    if st.button("Register"):
        if register(new_username, new_password):
            st.success("Registered successfully")
            st.session_state.page = 'login'
        else:
            st.error("Username already exists")

    if st.button("Go to Login"):
        st.session_state.page = 'login'

def show_collection_management_page():
    st.title("Collection Management")
    st.write(f"Logged in as: {st.session_state.current_user}")

    if st.button("Create New Collection"):
        st.session_state.page = 'create_collection'

    if st.button("Select Existing Collection"):
        st.session_state.page = 'select_collection'

    if st.button("Logout"):
        st.session_state.current_user = None
        st.session_state.page = 'login'

def show_create_collection_page():
    st.title("Create a New Collection")
    collection_name = st.text_input("Collection Name")
    index_name = st.text_input("Pinecone Index Name")  # User inputs the index name here
    saved_index_name = index_name
    st.write(f"Debug: index_name = '{index_name}'")  # Debugging print

    uploaded_files = st.file_uploader("Upload files", type=["txt", "pdf", "ppt", "csv"], accept_multiple_files=True)

    if st.button("Create"):
        if not saved_index_name or len(saved_index_name.strip()) == 0:  # Validate that index name is not empty
            st.error("Pinecone Index Name cannot be empty. Please provide a valid index name.")
        else:
            # Call process_and_upsert_files with saved_index_name
            result = process_and_upsert_files(st.session_state.current_user, collection_name, saved_index_name, uploaded_files)
            if result == "Success":
                st.success(f"Collection '{collection_name}' with index '{saved_index_name}' has been successfully created and files have been upserted.")
                add_collection(st.session_state.current_user, collection_name, saved_index_name, uploaded_files)  # Ensure add_collection is called after success
                st.session_state.page = 'collection_management'
            else:
                st.error(result)

    if st.button("Back to Collection Management"):
        st.session_state.page = 'collection_management'

def show_select_collection_page():
    st.title("Select an Existing Collection")

    # Update collection list for the logged-in user
    user_collections = st.session_state.collections.get(st.session_state.current_user, [])
    collection_names = [collection['name'] for collection in user_collections]
    selected_collection = st.selectbox("Select a collection", options=collection_names)

    if st.button("Select"):
        if selected_collection:
            st.session_state.selected_collection = selected_collection
            st.session_state.page = 'ask_questions'

    if st.button("Back to Collection Management"):
        st.session_state.page = 'collection_management'

def show_ask_questions_page():
    st.title("Ask a Question")
    st.write(f"Logged in as: {st.session_state.current_user}")
    st.write(f"Selected Collection: {st.session_state.selected_collection}")

    question = st.text_input("Ask a question")
    if st.button("Submit"):
        if question:
            st.success(f"Question submitted: {question}")
            # Add logic to handle the question, e.g., querying Pinecone for the answer

    if st.button("Back to Collection Management"):
        st.session_state.page = 'collection_management'

# Page Routing
if st.session_state.page == 'login':
    show_login_page()
elif st.session_state.page == 'register':
    show_register_page()
elif st.session_state.page == 'collection_management':
    show_collection_management_page()
elif st.session_state.page == 'create_collection':
    show_create_collection_page()
elif st.session_state.page == 'select_collection':
    show_select_collection_page()
elif st.session_state.page == 'ask_questions':
    show_ask_questions_page()
